"""One slide per page, and the slide count that proves it."""

import os
import zipfile

import pikepdf
import pytest
from pikepdf import Array, Dictionary, Name
from pptx import Presentation

from engine.office_export import export_document
from engine.slide_export import EMU_PER_POINT, _family


def _write(path, pages, size=(612, 792), rotate=0):
    pdf = pikepdf.new()
    for ops in pages:
        font = pdf.make_indirect(Dictionary(
            Type=Name.Font, Subtype=Name.Type1,
            BaseFont=Name.Helvetica, Encoding=Name.WinAnsiEncoding))
        bold = pdf.make_indirect(Dictionary(
            Type=Name.Font, Subtype=Name.Type1,
            BaseFont=Name("/Helvetica-Bold"), Encoding=Name.WinAnsiEncoding))
        page = Dictionary(
            Type=Name.Page, MediaBox=Array([0, 0, size[0], size[1]]),
            Resources=Dictionary(Font=Dictionary(F1=font, F2=bold)),
            Contents=pdf.make_stream("\n".join(ops).encode("latin-1")))
        if rotate:
            page[Name.Rotate] = rotate
        pdf.pages.append(pikepdf.Page(pdf.make_indirect(page)))
    pdf.save(str(path))
    pdf.close()
    return str(path)


def _text_page(n):
    return [
        f"BT /F1 18 Tf 72 720 Td (Heading of page {n}) Tj ET",
        f"BT /F1 11 Tf 72 690 Td (Body sentence on page {n}.) Tj ET",
    ]


def _slide_members(path):
    with zipfile.ZipFile(path) as parts:
        return [
            name for name in parts.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        ]


def _runs(path):
    out = []
    for slide in Presentation(path).slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                texts.append(shape.text_frame.text)
        out.append(texts)
    return out


def test_a_three_page_document_writes_exactly_three_slides(tmp_dir, gs_path):
    src = _write(os.path.join(tmp_dir, "s.pdf"), [_text_page(n) for n in (1, 2, 3)])
    out = os.path.join(tmp_dir, "s.pptx")
    result = export_document(src, out, "pptx", gs_path=gs_path)
    assert result["slides"] == 3
    assert result["pages_exported"] == [1, 2, 3]
    assert len(_slide_members(out)) == 3
    texts = _runs(out)
    assert "Heading of page 2" in texts[1]
    assert "Body sentence on page 3." in texts[2]


def test_the_slide_size_is_the_pages_own_size(tmp_dir, gs_path):
    src = _write(os.path.join(tmp_dir, "a5.pdf"), [_text_page(1)], size=(420, 595))
    out = os.path.join(tmp_dir, "a5.pptx")
    export_document(src, out, "pptx", gs_path=gs_path)
    deck = Presentation(out)
    assert round(deck.slide_width / EMU_PER_POINT) == 420
    assert round(deck.slide_height / EMU_PER_POINT) == 595


def test_a_preset_overrides_the_page_size(tmp_dir, gs_path):
    src = _write(os.path.join(tmp_dir, "s.pdf"), [_text_page(1)])
    out = os.path.join(tmp_dir, "s.pptx")
    result = export_document(src, out, "pptx", gs_path=gs_path, slide_size="4:3")
    assert result["slide_size"] == "4:3"
    deck = Presentation(out)
    assert round(deck.slide_width / EMU_PER_POINT) == 720
    assert round(deck.slide_height / EMU_PER_POINT) == 540


def test_pages_of_a_different_size_are_counted_not_silently_rescaled(tmp_dir, gs_path):
    src = _write(os.path.join(tmp_dir, "one.pdf"), [_text_page(1)])
    with pikepdf.open(src, allow_overwriting_input=True) as pdf:
        second = pikepdf.Page(pdf.make_indirect(Dictionary(
            Type=Name.Page, MediaBox=Array([0, 0, 595, 842]),
            Resources=Dictionary(),
            Contents=pdf.make_stream(b""))))
        pdf.pages.append(second)
        pdf.save(src)
    result = export_document(src, os.path.join(tmp_dir, "one.pptx"), "pptx", gs_path=gs_path)
    assert result["slides"] == 2
    assert result["pages_of_a_different_size"] == 1


def test_a_rotated_page_carries_its_text_at_the_display_position(tmp_dir, gs_path):
    ops = _text_page(1)
    src = _write(os.path.join(tmp_dir, "turned.pdf"), [ops], rotate=90)
    out = os.path.join(tmp_dir, "turned.pptx")
    export_document(src, out, "pptx", gs_path=gs_path)
    deck = Presentation(out)
    # The deck takes the DISPLAY dimensions, so a quarter-turned Letter page is
    # a landscape slide.
    assert round(deck.slide_width / EMU_PER_POINT) == 792
    assert round(deck.slide_height / EMU_PER_POINT) == 612
    shapes = [s for s in deck.slides[0].shapes if s.has_text_frame and s.text_frame.text]
    assert {s.text_frame.text for s in shapes} == {
        "Heading of page 1", "Body sentence on page 1."
    }
    assert all(s.rotation == 90 for s in shapes)


def test_the_page_graphics_land_on_the_slide(tmp_dir, gs_path):
    ops = ["1 0 0 RG 4 w 100 400 m 500 400 l S"] + _text_page(1)
    src = _write(os.path.join(tmp_dir, "drawn.pdf"), [ops])
    out = os.path.join(tmp_dir, "drawn.pptx")
    result = export_document(src, out, "pptx", gs_path=gs_path)
    assert result["rasterized_pages"] == 1
    assert any(s.shape_type == 13 for s in Presentation(out).slides[0].shapes)


def test_a_page_with_nothing_but_text_carries_no_raster(tmp_dir, gs_path):
    # The background is rendered WITHOUT text, so a text-only page renders blank
    # and placing that image would only add weight.
    src = _write(os.path.join(tmp_dir, "plain.pdf"), [_text_page(1)])
    out = os.path.join(tmp_dir, "plain.pptx")
    result = export_document(src, out, "pptx", gs_path=gs_path)
    assert result["rasterized_pages"] == 0
    assert result["text_boxes"] == 2


def test_a_deck_with_no_slides_is_refused_and_nothing_survives(tmp_dir, gs_path, monkeypatch):
    # A presentation package's success is its SLIDE COUNT: a writer can produce
    # a well-formed, non-empty package carrying no slides at all, and a size
    # check passes it.
    import engine.slide_export as se

    src = _write(os.path.join(tmp_dir, "s.pdf"), [_text_page(1), _text_page(2)])
    out = os.path.join(tmp_dir, "s.pptx")
    monkeypatch.setattr(se, "_count_slides", lambda _package: 0)
    with pytest.raises(RuntimeError, match="0 slide"):
        export_document(src, out, "pptx", gs_path=gs_path)
    assert not os.path.exists(out)


def test_refuses_without_ghostscript(tmp_dir):
    src = _write(os.path.join(tmp_dir, "s.pdf"), [_text_page(1)])
    with pytest.raises(RuntimeError, match="Ghostscript is not available"):
        export_document(src, os.path.join(tmp_dir, "s.pptx"), "pptx", gs_path="")


def test_rejects_an_unknown_slide_size(tmp_dir, gs_path):
    src = _write(os.path.join(tmp_dir, "s.pdf"), [_text_page(1)])
    with pytest.raises(ValueError, match="unknown slide size"):
        export_document(
            src, os.path.join(tmp_dir, "s.pptx"), "pptx", gs_path=gs_path, slide_size="A4"
        )


def test_the_option_belongs_to_the_presentation_target_only(tmp_dir, gs_path):
    src = _write(os.path.join(tmp_dir, "s.pdf"), [_text_page(1)])
    with pytest.raises(ValueError, match="the xlsx export takes no slide_size option"):
        export_document(src, os.path.join(tmp_dir, "s.xlsx"), "xlsx", slide_size="4:3")


def test_a_page_scope_limits_the_deck(tmp_dir, gs_path):
    src = _write(os.path.join(tmp_dir, "s.pdf"), [_text_page(n) for n in (1, 2, 3)])
    out = os.path.join(tmp_dir, "s.pptx")
    result = export_document(src, out, "pptx", gs_path=gs_path, pages=[2])
    assert result["slides"] == 1
    assert result["pages_exported"] == [2]
    assert _runs(out)[0] == ["Heading of page 2", "Body sentence on page 2."]


def test_the_face_a_run_drew_with_reaches_the_slide(tmp_dir, gs_path):
    ops = ["BT /F2 18 Tf 72 720 Td (A bold heading) Tj ET"]
    src = _write(os.path.join(tmp_dir, "bold.pdf"), [ops])
    out = os.path.join(tmp_dir, "bold.pptx")
    export_document(src, out, "pptx", gs_path=gs_path)
    shape = next(
        s for s in Presentation(out).slides[0].shapes
        if s.has_text_frame and s.text_frame.text
    )
    run = shape.text_frame.paragraphs[0].runs[0]
    assert run.font.bold is True
    assert run.font.name == "Arial"
    assert round(run.font.size / EMU_PER_POINT) == 18


@pytest.mark.parametrize(
    "base,family,bold,italic",
    [
        ("/ABCDEF+Helvetica-Bold", "Arial", True, False),
        ("/TimesNewRomanPSMT", "Times New Roman", False, False),
        ("/Calibri-Italic", "Calibri", False, True),
        ("/CourierNewPS-BoldItalicMT", "Courier New", True, True),
    ],
)
def test_a_base_font_resolves_to_a_family_and_a_style(base, family, bold, italic):
    assert _family(base) == (family, bold, italic)
