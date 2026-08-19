"""Turn a PDF's pages into presentation slides, one slide per page.

A presentation is a layout transcription, so a slide carries the page twice
over, in the two forms a deck needs:

* **Text as text**, in boxes at their measured positions, at the size and face
  the page drew them with. A deck of page pictures is not editable and is not
  what the target is for.
* **Everything else as a raster**, rendered without its text so the two layers
  never draw the same glyph twice. A deck that dropped the page's rules,
  images and fills would be half a page.

Success is the SLIDE COUNT, read back off the written package. A presentation
package's success is its slide count: a conversion can write a well-formed
package with a non-zero size and no slides in it at all.
"""

from __future__ import annotations

import re
import shutil
import tempfile
import zipfile
from pathlib import Path

import pikepdf

from engine import bidi, budget
from engine.form_detect import _crop_box, _page_rotate, _page_segments
from engine.redact import _resolve_resources
from engine.soffice import _normalise_face

# What the background raster is rendered at. Matches the image export's own
# default: high enough that a rule or a logo survives, low enough that a long
# document is not gigabytes of PNG.
RASTER_DPI = 150
# One point in English Metric Units, the unit a presentation package measures in.
EMU_PER_POINT = 12700
# The largest slide edge the format admits, in points.
MAX_SLIDE_POINTS = 56 * 72
SLIDE_SIZES = ("page", "16:9", "4:3")
_PRESET_POINTS = {"16:9": (13.333 * 72, 7.5 * 72), "4:3": (10.0 * 72, 7.5 * 72)}

_SUBSET_TAG = re.compile(r"^[A-Z]{6}\+")
# The families a document names through their PostScript aliases. Anything else
# keeps the name the document gave, which is the name a reader's machine has the
# best chance of matching.
_STANDARD_FAMILIES = {
    "helvetica": "Arial",
    "arial": "Arial",
    "times": "Times New Roman",
    "timesnewroman": "Times New Roman",
    "courier": "Courier New",
    "couriernew": "Courier New",
    "symbol": "Symbol",
}
FALLBACK_FAMILY = "Arial"
_BOLD = re.compile(r"bold|black|heavy|semibold", re.I)
_ITALIC = re.compile(r"italic|oblique", re.I)


def _family(base_font: str) -> tuple[str, bool, bool]:
    """(family name, bold, italic) from an embedded font's base name."""
    name = _SUBSET_TAG.sub("", str(base_font or "").lstrip("/"))
    style = name
    name = re.split(r"[-,]", name, maxsplit=1)[0]
    key = _normalise_face(name)
    family = _STANDARD_FAMILIES.get(key, name or FALLBACK_FAMILY)
    return family, bool(_BOLD.search(style)), bool(_ITALIC.search(style))


def _base_fonts(page) -> dict:
    """Resource name -> base font name, for the fonts this page selects."""
    out: dict[str, str] = {}
    try:
        fonts = (_resolve_resources(page) or {}).get("/Font") or {}
    except (AttributeError, TypeError, KeyError):
        return out
    for key, font in fonts.items():
        try:
            base = font.get("/BaseFont")
            if base is None:
                descendants = font.get("/DescendantFonts") or []
                base = descendants[0].get("/BaseFont") if len(descendants) else None
            if base is not None:
                out[str(key)] = str(base)
        except (AttributeError, TypeError, IndexError):
            continue
    return out


def _render_background(file: str, page_number: int, gs_path: str, out_png: Path) -> None:
    """Rasterise one page WITHOUT its text.

    The text is drawn separately as editable boxes, so a raster that kept it
    would double every glyph — visibly, because a box's own line breaking never
    lands exactly where the page's did.

    -dUseCropBox frames the raster on the box the slide is built from: the
    picture is placed at `_display_size(_crop_box(page), rotate)` and every text
    box is positioned against that same box. A MediaBox raster of a cropped page
    is squeezed into the crop-sized picture frame, so the graphics land at a
    different scale from the text drawn over them. The device clips page content
    to the CropBox either way, and a page with no CropBox is unaffected.
    """
    cmd = [
        gs_path,
        "-q",
        "-dNOPAUSE",
        "-dBATCH",
        "-dSAFER",
        "-sDEVICE=png16m",
        "-dFILTERTEXT",
        "-dTextAlphaBits=4",
        "-dGraphicsAlphaBits=4",
        f"-r{RASTER_DPI}",
        "-dUseCropBox",
        f"-dFirstPage={page_number}",
        f"-dLastPage={page_number}",
        f"-sOutputFile={out_png}",
        str(file),
    ]
    # Through the gs family's door: it validates the executable (an existing
    # file is not a working interpreter) and derives the budget. This run had
    # no timeout before, so a wedged interpreter stalled the export with
    # nothing to report.
    proc = budget.gs(cmd, what=f"Ghostscript (slide raster, page {page_number})",
                     path=file, pages=1)
    if proc.returncode != 0 or not out_png.is_file():
        detail = (proc.stderr or proc.stdout or "").strip()[:400]
        raise RuntimeError(f"Could not render page {page_number} for a slide: {detail}")


def _is_blank(png: Path) -> bool:
    """A render of one colour carries nothing; placing it would only add weight."""
    from PIL import Image

    with Image.open(png) as image:
        # getcolors returns None above its limit, so a limit of one answers
        # "exactly one colour" — a limit of two would call a solid logo on white
        # blank and drop it.
        return image.convert("RGB").getcolors(1) is not None


def _display_point(x: float, y: float, page_box: tuple, rotate: int) -> tuple:
    """A user-space point in DISPLAY space, origin top-left.

    The rasteriser applies the page's `/Rotate`, so the background image is in
    display orientation while every rect the text channel reports is in
    un-rotated user space.
    """
    mx, my = page_box[0], page_box[1]
    width, height = page_box[2] - page_box[0], page_box[3] - page_box[1]
    u, v = x - mx, y - my
    angle = int(rotate) % 360
    if angle == 90:
        return (v, u)
    if angle == 180:
        return (width - u, v)
    if angle == 270:
        return (height - v, width - u)
    return (u, height - v)


def _display_size(page_box: tuple, rotate: int) -> tuple:
    width, height = page_box[2] - page_box[0], page_box[3] - page_box[1]
    return (height, width) if int(rotate) % 360 in (90, 270) else (width, height)


def _logical(text: str) -> tuple[str, bool]:
    """(the run's text in logical order, whether it carries right-to-left text).

    A run is assembled left to right by geometry, so page order is VISUAL order.
    The inverse is proven by re-running the forward reordering and requiring the
    permutation to compose to the identity; an unprovable run is kept as drawn.
    """
    if not bidi.has_strong_rtl(text):
        return text, False
    back = bidi.reconstruct_logical(text, 1)
    if len(back) != len(text):
        return text, True
    logical = "".join(text[i] for i in back)
    _level, forward = bidi.visual_order(logical, 1)
    if len(forward) != len(text) or any(back[forward[v]] != v for v in range(len(text))):
        return text, True
    return logical, True


def _page_numbers(pages, pdf) -> list[int]:
    """The 1-based pages to export.

    Spelled with `page_no` and `len(pdf.pages)` deliberately: the out-of-range
    refusal is a shared row in the engine-message table, and a differently named
    local would rename the interpolations of a message several modules raise.
    """
    if pages is None or pages == "all":
        return list(range(1, len(pdf.pages) + 1))
    if isinstance(pages, str):
        raise ValueError('pages must be a list of page numbers or "all"')
    out: list[int] = []
    for value in pages:
        page_no = int(value)
        if not (1 <= page_no <= len(pdf.pages)):
            raise ValueError(f"page {page_no} is out of range (1-{len(pdf.pages)})")
        if page_no not in out:
            out.append(page_no)
    return sorted(out)


def _slide_points(sizing: str, first: tuple) -> tuple:
    if sizing in _PRESET_POINTS:
        return _PRESET_POINTS[sizing]
    width, height = first
    scale = min(1.0, MAX_SLIDE_POINTS / max(width, height, 1.0))
    return (width * scale, height * scale)


def _add_text(slide, segment, base_fonts, page_box, rotate, scale, offset) -> None:
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Emu, Pt

    x0, y0, x1, y1 = segment.rect
    width = max(x1 - x0, 1.0) * scale
    height = max(y1 - y0, 1.0) * scale
    centre = _display_point((x0 + x1) / 2.0, (y0 + y1) / 2.0, page_box, rotate)
    left = offset[0] + centre[0] * scale - width / 2.0
    top = offset[1] + centre[1] * scale - height / 2.0

    box = slide.shapes.add_textbox(
        Emu(int(left * EMU_PER_POINT)),
        Emu(int(top * EMU_PER_POINT)),
        Emu(int(width * EMU_PER_POINT)),
        Emu(int(height * EMU_PER_POINT)),
    )
    if int(rotate) % 360:
        # The box rotates about its own centre, which is why the centre is what
        # was mapped: a rotated page draws its text along the rotated axis.
        box.rotation = int(rotate) % 360
    frame = box.text_frame
    frame.word_wrap = False
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    text, is_rtl = _logical(segment.text)
    if is_rtl:
        # The package's own right-to-left paragraph flag; the text itself is
        # already in logical order, which is what the flag describes.
        paragraph._p.get_or_add_pPr().set("rtl", "1")
    run = paragraph.add_run()
    run.text = text
    family, bold, italic = _family(base_fonts.get(segment.font_name, ""))
    run.font.size = Pt(max(segment.size, 1.0) * scale)
    run.font.name = family
    run.font.bold = bold
    run.font.italic = italic


def export_slides(
    file: str,
    output: str,
    pages="all",
    slide_size: str = "page",
    gs_path: str = "",
) -> dict:
    """Write ``file``'s pages to the presentation at ``output``, one slide each.

    Args:
        file: input PDF path.
        output: destination ``.pptx`` path.
        pages: list of 1-based page numbers, or 'all'.
        slide_size: 'page' (the document's own size) or a '16:9' / '4:3' preset.
        gs_path: path to the Ghostscript executable, for the page graphics.
    """
    from pptx import Presentation
    from pptx.util import Emu

    sizing = str(slide_size or "page").lower()
    if sizing not in SLIDE_SIZES:
        raise ValueError(
            f"unknown slide size {slide_size!r} (choose page, 16:9 or 4:3)"
        )

    work = Path(tempfile.mkdtemp(prefix="slide-export-"))
    try:
        with pikepdf.open(str(file)) as pdf:
            wanted = _page_numbers(pages, pdf)
            geometry = []
            for number in wanted:
                page = pdf.pages[number - 1]
                box = _crop_box(page)
                rotate = _page_rotate(page)
                geometry.append((number, page, box, rotate, _display_size(box, rotate)))

            first = geometry[0][4]
            differing = sum(1 for item in geometry if item[4] != first)
            deck_width, deck_height = _slide_points(sizing, first)

            deck = Presentation()
            deck.slide_width = Emu(int(deck_width * EMU_PER_POINT))
            deck.slide_height = Emu(int(deck_height * EMU_PER_POINT))
            blank = deck.slide_layouts[6]

            text_boxes = 0
            rasterized = 0
            for number, page, box, rotate, size in geometry:
                slide = deck.slides.add_slide(blank)
                # A page whose size differs from the deck's is fitted rather
                # than cropped: content placed past the slide edge is content
                # nobody can see, and the count of such pages is reported.
                scale = min(deck_width / size[0], deck_height / size[1])
                offset = (
                    (deck_width - size[0] * scale) / 2.0,
                    (deck_height - size[1] * scale) / 2.0,
                )
                png = work / f"page{number}.png"
                _render_background(str(file), number, gs_path, png)
                if not _is_blank(png):
                    slide.shapes.add_picture(
                        str(png),
                        Emu(int(offset[0] * EMU_PER_POINT)),
                        Emu(int(offset[1] * EMU_PER_POINT)),
                        width=Emu(int(size[0] * scale * EMU_PER_POINT)),
                        height=Emu(int(size[1] * scale * EMU_PER_POINT)),
                    )
                    rasterized += 1
                base_fonts = _base_fonts(page)
                for segment in _page_segments(pdf, page):
                    _add_text(slide, segment, base_fonts, box, rotate, scale, offset)
                    text_boxes += 1

        out_path = Path(output)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        staged = work / "deck.pptx"
        deck.save(str(staged))

        slides = _count_slides(staged)
        if slides != len(wanted):
            raise RuntimeError(
                f"the presentation was written with {slides} slide(s) for "
                f"{len(wanted)} exported page(s), so it does not carry the document"
            )
        if out_path.exists():
            try:
                out_path.chmod(0o666)
            except OSError:
                pass
        shutil.move(str(staged), str(out_path))
        return {
            "output": str(out_path),
            "format": "pptx",
            "size": out_path.stat().st_size,
            "slides": slides,
            "pages_exported": wanted,
            "text_boxes": text_boxes,
            "rasterized_pages": rasterized,
            "pages_of_a_different_size": differing,
            "slide_size": sizing,
        }
    finally:
        shutil.rmtree(work, ignore_errors=True)


def _count_slides(package: Path) -> int:
    """How many slides the written package actually carries."""
    with zipfile.ZipFile(package) as parts:
        return sum(
            1
            for name in parts.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
